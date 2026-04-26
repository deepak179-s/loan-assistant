# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

# ─── PALETTE ────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x14, 0x2B)   # deep navy background
BG_MID     = RGBColor(0x13, 0x1E, 0x3A)   # slightly lighter navy
ACCENT1    = RGBColor(0x00, 0xC9, 0xA7)   # teal/mint
ACCENT2    = RGBColor(0xFF, 0x6B, 0x35)   # orange
ACCENT3    = RGBColor(0x7C, 0x3A, 0xFF)   # purple
CARD_BG    = RGBColor(0x1A, 0x28, 0x4D)   # card background
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x8A, 0x9B, 0xC4)   # muted blue-grey text
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)
RED_ALERT  = RGBColor(0xFF, 0x3B, 0x3B)
GREEN_OK   = RGBColor(0x00, 0xE0, 0x7A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ─── HELPERS ────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0), alpha=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, font="Trebuchet MS"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def multiline_txt(slide, lines, x, y, w, h, size=14, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, font="Trebuchet MS", line_spacing=None):
    """lines = list of (text, bold, color, size) tuples or plain strings"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            t, b, c, s = item, bold, color, size
        else:
            t, b, c, s = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            from pptx.util import Pt as _Pt
            p.space_after = _Pt(line_spacing)
        run = p.add_run()
        run.text = t
        run.font.name = font
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
    return txb

def card(slide, x, y, w, h, fill=CARD_BG, border=ACCENT1, bw=Pt(1.2)):
    r = rect(slide, x, y, w, h, fill_rgb=fill, line_rgb=border, line_w=bw)
    return r

def accent_bar(slide, x, y, h=0.04, w=1.0, color=ACCENT1):
    rect(slide, x, y, w, h, fill_rgb=color)

def slide_number(slide, n, total=12):
    txt(slide, f"{n} / {total}", 12.3, 7.1, 0.8, 0.3,
        size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def tag_box(slide, label, x, y, fill=ACCENT1):
    r = rect(slide, x, y, 1.5, 0.28, fill_rgb=fill)
    txt(slide, label, x+0.05, y+0.03, 1.4, 0.22, size=9, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER)

def image_placeholder(slide, x, y, w, h, label="[ INSERT IMAGE / DIAGRAM HERE ]"):
    """Dashed border placeholder for images."""
    from pptx.util import Pt
    # outer box
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0A, 0x12, 0x26)
    shape.line.color.rgb = ACCENT1
    shape.line.width = Pt(1.2)
    shape.line.dash_style = 4   # DASH
    # label in center
    txb = slide.shapes.add_textbox(
        Inches(x+0.1), Inches(y + h/2 - 0.25), Inches(w-0.2), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Trebuchet MS"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ACCENT1
    run.font.italic = True
    return shape

# ════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)

# left dark panel
rect(s1, 0, 0, 6.3, 7.5, fill_rgb=BG_MID)

# teal vertical accent strip
rect(s1, 0, 0, 0.12, 7.5, fill_rgb=ACCENT1)

# tag
tag_box(s1, "INT428 · IOP", 0.4, 0.55, fill=ACCENT1)

# main title
multiline_txt(s1, [
    ("AI-Powered", True, WHITE, 36),
    ("Financial Debt", True, WHITE, 36),
    ("Optimizer", True, ACCENT1, 40),
], 0.4, 1.1, 5.7, 2.2, font="Trebuchet MS")

txt(s1, "& Loan Assistant Platform", 0.4, 3.35, 5.5, 0.5,
    size=18, bold=False, color=MUTED)

# divider
rect(s1, 0.4, 3.95, 3.0, 0.04, fill_rgb=ACCENT1)

# presenter info
multiline_txt(s1, [
    ("Presented By: Deepak Kumar", True, WHITE, 13),
    ("UID: 22BCS14872", False, MUTED, 11),
    ("B.E. — Computer Science & Engineering", False, MUTED, 11),
    ("Chandigarh University", False, MUTED, 11),
], 0.4, 4.15, 5.5, 1.2, font="Trebuchet MS")

# tech stack labels
txt(s1, "TECH STACK", 0.4, 5.55, 5.5, 0.28,
    size=9, bold=True, color=MUTED, font="Trebuchet MS")
tech = ["React.js", "Node.js", "Firebase", "Groq · Llama-3 70B"]
cols = [0.4, 1.85, 3.25, 4.5]
for i, t in enumerate(tech):
    rect(s1, cols[i], 5.9, 1.25, 0.38,
         fill_rgb=RGBColor(0x1A,0x28,0x4D), line_rgb=ACCENT1, line_w=Pt(0.7))
    txt(s1, t, cols[i]+0.07, 5.97, 1.1, 0.24,
        size=9, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

# right side: image placeholder
image_placeholder(s1, 6.6, 0.4, 6.4, 6.7,
    "[ Screenshot: Dashboard with\n  CIBIL card + Area Chart ]")

slide_number(s1, 1)

# ════════════════════════════════════════════════════════════════
#  SLIDE 2 — CORE PROBLEMS
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)

# top accent strip
rect(s2, 0, 0, 13.33, 0.07, fill_rgb=ACCENT2)

tag_box(s2, "THE PROBLEM", 0.5, 0.22, fill=ACCENT2)
txt(s2, "Why Current FinTech AI Fails", 0.5, 0.6, 9.0, 0.7,
    size=30, bold=True, color=WHITE, font="Trebuchet MS")
txt(s2, "Three critical, unsolved flaws in every existing AI financial chatbot",
    0.5, 1.3, 9.0, 0.38, size=13, color=MUTED, font="Trebuchet MS")

# 3 problem cards
problems = [
    ("01", "Mathematical\nHallucinations", ACCENT2,
     "LLMs are text predictors — not calculators.\nThey fabricate amortization results, producing\nfigures that could cost users crores of rupees."),
    ("02", "Fraud\nSusceptibility", RED_ALERT,
     "Standard chatbots cannot distinguish between\na valid strategy and a WhatsApp lottery scam.\nNo guardrails exist in generic AI wrappers."),
    ("03", "Data Privacy\n(DPDP Violation)", ACCENT3,
     "Asking users to type PAN numbers and loan\nbalances into chat windows violates India's\nDigital Personal Data Protection Act, 2023."),
]
card_x = [0.5, 4.6, 8.7]
for i, (num, title, col, desc) in enumerate(problems):
    cx = card_x[i]
    card(s2, cx, 1.85, 4.0, 4.7, fill=CARD_BG, border=col, bw=Pt(1.5))
    # number badge
    rect(s2, cx+0.18, 2.0, 0.52, 0.52, fill_rgb=col)
    txt(s2, num, cx+0.18, 2.0, 0.52, 0.52, size=13, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER, font="Trebuchet MS")
    # title
    txt(s2, title, cx+0.85, 1.95, 3.0, 0.65,
        size=15, bold=True, color=col, font="Trebuchet MS")
    # desc
    txt(s2, desc, cx+0.18, 2.7, 3.65, 2.5,
        size=12, color=MUTED, font="Trebuchet MS")

# right image placeholder
image_placeholder(s2, 8.7, 1.85, 4.0, 4.7,
    "[ Side-by-side: ChatGPT\n  math fail + Scam WhatsApp ]")

slide_number(s2, 2)

# ════════════════════════════════════════════════════════════════
#  SLIDE 3 — SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
rect(s3, 0, 0, 13.33, 0.07, fill_rgb=ACCENT1)

tag_box(s3, "ARCHITECTURE", 0.5, 0.22, fill=ACCENT1)
txt(s3, "System Architecture & Tech Stack", 0.5, 0.6, 9.0, 0.65,
    size=28, bold=True, color=WHITE, font="Trebuchet MS")
txt(s3, "Decoupled 3-tier design — security by isolation",
    0.5, 1.25, 7.0, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# left col: tech stack cards
layers = [
    ("Frontend",    "React.js 19 · Vite · TypeScript\nRecharts · Lucide-React · React-Router", ACCENT1),
    ("Backend",     "Node.js · Express.js 5\nVercel Serverless Functions",                     ACCENT2),
    ("Database",    "Firebase Firestore · Firebase Auth\nOTP-based KYC Verification",          ACCENT3),
    ("AI Engine",   "Groq Cloud · Llama-3.3-70B-Versatile\nDual Agent Inference Pipeline",    YELLOW),
]
for i, (layer, detail, col) in enumerate(layers):
    cy = 1.75 + i * 1.2
    card(s3, 0.5, cy, 5.8, 1.05, fill=CARD_BG, border=col, bw=Pt(1.2))
    rect(s3, 0.5, cy, 0.08, 1.05, fill_rgb=col)
    txt(s3, layer, 0.75, cy+0.1, 1.5, 0.35, size=12, bold=True, color=col,
        font="Trebuchet MS")
    txt(s3, detail, 0.75, cy+0.45, 5.3, 0.55, size=11, color=MUTED,
        font="Trebuchet MS")

# right: architecture diagram placeholder
image_placeholder(s3, 6.7, 1.55, 6.3, 5.5,
    "[ FIGURE 1:\n  System Architecture Diagram\n  (Frontend ↔ Backend ↔ Firebase / Groq) ]")

slide_number(s3, 3)

# ════════════════════════════════════════════════════════════════
#  SLIDE 4 — MULTI-AGENT ARCHITECTURE
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
rect(s4, 0, 0, 13.33, 0.07, fill_rgb=ACCENT3)

tag_box(s4, "CORE INNOVATION", 0.5, 0.22, fill=ACCENT3)
txt(s4, "Multi-Agent Debating Architecture", 0.5, 0.6, 9.0, 0.65,
    size=28, bold=True, color=WHITE, font="Trebuchet MS")
txt(s4, "Two adversarial LLM agents work in sequence — one proposes, one audits",
    0.5, 1.25, 9.5, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# agent cards - left column
agents = [
    ("Agent 1", "Lead Strategist", ACCENT1,
     "• Receives user query + injected JSON profile\n"
     "• Proposes Snowball / Avalanche debt strategy\n"
     "• Explicitly BLOCKED from computing exact math"),
    ("Agent 2", "Expert Critic", ACCENT3,
     "• Independently audits the Strategist's draft\n"
     "• Scans for impossible math & fraud signals\n"
     "• Outputs JSON: { confidenceScore, finalAdvice }"),
]
for i, (tag, title, col, desc) in enumerate(agents):
    cy = 1.8 + i * 2.1
    card(s4, 0.5, cy, 5.8, 1.85, fill=CARD_BG, border=col, bw=Pt(1.5))
    tag_box(s4, tag, 0.65, cy+0.18, fill=col)
    txt(s4, title, 2.3, cy+0.15, 4.0, 0.38,
        size=15, bold=True, color=col, font="Trebuchet MS")
    txt(s4, desc, 0.65, cy+0.65, 5.5, 1.1,
        size=12, color=MUTED, font="Trebuchet MS")

# JSON output box
card(s4, 0.5, 6.05, 5.8, 0.9, fill=RGBColor(0x0D,0x0D,0x2B), border=YELLOW, bw=Pt(1))
txt(s4, "OUTPUT  →  { confidenceScore: 85, reasoning: '...', finalAdvice: '...' }",
    0.65, 6.18, 5.55, 0.62, size=11, bold=True, color=YELLOW, font="Consolas")

# right: flow diagram placeholder
image_placeholder(s4, 6.7, 1.6, 6.3, 5.5,
    "[ FIGURE 2:\n  Multi-Agent Debating Loop\n  Flow Diagram ]")

slide_number(s4, 4)

# ════════════════════════════════════════════════════════════════
#  SLIDE 5 — UI SCAM INTERCEPTION
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
rect(s5, 0, 0, 13.33, 0.07, fill_rgb=RED_ALERT)

tag_box(s5, "SAFETY LAYER", 0.5, 0.22, fill=RED_ALERT)
txt(s5, "UI Scam Interception & Human Handoff", 0.5, 0.6, 9.0, 0.65,
    size=28, bold=True, color=WHITE, font="Trebuchet MS")
txt(s5, "React intercepts the confidence score and dynamically routes rendering",
    0.5, 1.25, 9.5, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# 3 routing tiers
tiers = [
    ("Score 80 – 100", "SAFE",     GREEN_OK,  "Normal AI advice rendered in chat.\nLoRA dataset entry logged automatically."),
    ("Score 31 – 79",  "CAUTION",  YELLOW,    "Orange warning banner shown below advice.\nHuman Advisor link provided to user."),
    ("Score 0 – 30",   "BLOCKED",  RED_ALERT, "Red  🚨  Scam Alert module replaces chat.\nAdvice suppressed. Mandatory Human Callback."),
]
for i, (score, label, col, desc) in enumerate(tiers):
    cy = 1.8 + i * 1.55
    card(s5, 0.5, cy, 5.6, 1.3, fill=CARD_BG, border=col, bw=Pt(1.5))
    rect(s5, 0.5, cy, 0.08, 1.3, fill_rgb=col)
    txt(s5, score, 0.72, cy+0.12, 2.8, 0.3, size=11, bold=False, color=MUTED,
        font="Trebuchet MS")
    txt(s5, label, 0.72, cy+0.4, 2.5, 0.38, size=16, bold=True, color=col,
        font="Trebuchet MS")
    txt(s5, desc, 0.72, cy+0.78, 5.2, 0.48, size=11, color=MUTED,
        font="Trebuchet MS")

# right placeholder
image_placeholder(s5, 6.5, 1.6, 6.5, 5.5,
    "[ Screenshot: AI Agents chat showing\n  red  🚨  High Risk / Scam Detected\n  module in the UI ]")

# flowchart placeholder below
txt(s5, "← See Figure 4 for the complete Confidence Score Routing Flowchart",
    0.5, 6.95, 7.0, 0.38, size=10, color=MUTED, italic=True, font="Trebuchet MS")

slide_number(s5, 5)

# ════════════════════════════════════════════════════════════════
#  SLIDE 6 — DETERMINISTIC MATH VS AI
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_bg(s6)
rect(s6, 0, 0, 13.33, 0.07, fill_rgb=ACCENT1)

tag_box(s6, "INNOVATION", 0.5, 0.22, fill=ACCENT1)
txt(s6, "Deterministic Math vs. AI Reasoning", 0.5, 0.6, 9.0, 0.65,
    size=28, bold=True, color=WHITE, font="Trebuchet MS")
txt(s6, "Clear separation of concerns eliminates all mathematical hallucinations",
    0.5, 1.25, 9.5, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# two columns
# AI column
card(s6, 0.5, 1.75, 5.6, 4.9, fill=CARD_BG, border=ACCENT3, bw=Pt(1.2))
rect(s6, 0.5, 1.75, 5.6, 0.42, fill_rgb=ACCENT3)
txt(s6, "🤖  AI's Job — Strategic Reasoning", 0.65, 1.82, 5.3, 0.3,
    size=13, bold=True, color=WHITE, font="Trebuchet MS")
ai_points = [
    "Snowball vs. Avalanche prioritization",
    "Section 80E tax deduction strategy",
    "Balance transfer recommendations",
    "Macroeconomic context awareness",
    "Scam and fraud signal detection",
]
for j, pt in enumerate(ai_points):
    txt(s6, f"▸  {pt}", 0.65, 2.32 + j*0.6, 5.2, 0.48,
        size=12, color=MUTED, font="Trebuchet MS")

# React column
card(s6, 6.7, 1.75, 5.6, 4.9, fill=CARD_BG, border=ACCENT1, bw=Pt(1.2))
rect(s6, 6.7, 1.75, 5.6, 0.42, fill_rgb=ACCENT1)
txt(s6, "⚙️  React's Job — Exact Calculation", 6.85, 1.82, 5.3, 0.3,
    size=13, bold=True, color=BG_DARK, font="Trebuchet MS")
calc_points = [
    "Remaining EMIs  =  tenure − months_paid",
    "Principal Paid  =  original − outstanding",
    "Interest Paid  =  (EMI × months) − principal",
    "Annual amortization projection loop",
    "Credit utilization % per card",
]
for j, pt in enumerate(calc_points):
    txt(s6, f"▸  {pt}", 6.85, 2.32 + j*0.6, 5.3, 0.48,
        size=12, color=MUTED, font="Trebuchet MS")

# bottom image placeholder
image_placeholder(s6, 0.5, 6.75, 11.8, 0.6,
    "[ Screenshot: Credit Profile page — expanded loan card showing computed Principal Paid, Interest Paid, Remaining EMIs ]")

slide_number(s6, 6)

# ════════════════════════════════════════════════════════════════
#  SLIDE 7 — SECURE CONTEXT INJECTION
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_bg(s7)
rect(s7, 0, 0, 13.33, 0.07, fill_rgb=ACCENT2)

tag_box(s7, "DATA PRIVACY", 0.5, 0.22, fill=ACCENT2)
txt(s7, "Secure Context Injection  (DPDP Compliance)", 0.5, 0.6, 10.0, 0.65,
    size=27, bold=True, color=WHITE, font="Trebuchet MS")
txt(s7, "A proprietary Structured Data RAG pipeline — no raw financial data ever touches the frontend",
    0.5, 1.25, 12.0, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# left: pipeline diagram placeholder
image_placeholder(s7, 0.5, 1.7, 6.0, 5.5,
    "[ FIGURE 3:\n  Secure Context Injection\n  Pipeline Diagram ]")

# right: step-by-step
steps = [
    ("Step 1", "Frontend sends ONLY { userId, query }",          ACCENT2, "No financial data exposed in browser network."),
    ("Step 2", "Backend fetches creditProfile from Firestore",   ACCENT1, "Authenticated server-side read — user never sees this."),
    ("Step 3", "JSON serialized with [CONFIDENTIAL CONTEXT] tag",ACCENT3, "PAN, loan balances, rates injected into system prompt."),
    ("Step 4", "Groq API invoked with complete secure context",  YELLOW,  "LLM reasons over real data. Frontend remains isolated."),
]
for i, (step, title, col, sub) in enumerate(steps):
    cy = 1.7 + i * 1.35
    card(s7, 6.8, cy, 6.1, 1.18, fill=CARD_BG, border=col, bw=Pt(1))
    rect(s7, 6.8, cy, 0.07, 1.18, fill_rgb=col)
    txt(s7, step, 7.0, cy+0.1, 1.0, 0.28, size=9, bold=True, color=col,
        font="Trebuchet MS")
    txt(s7, title, 7.0, cy+0.35, 5.7, 0.35, size=12, bold=True, color=WHITE,
        font="Trebuchet MS")
    txt(s7, sub,   7.0, cy+0.72, 5.7, 0.38, size=11, color=MUTED,
        font="Trebuchet MS")

slide_number(s7, 7)

# ════════════════════════════════════════════════════════════════
#  SLIDE 8 — EMI SIMULATOR & XAI
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
rect(s8, 0, 0, 13.33, 0.07, fill_rgb=ACCENT1)

tag_box(s8, "SIMULATION", 0.5, 0.22, fill=ACCENT1)
txt(s8, "EMI Simulator & Explainable AI (SHAP)", 0.5, 0.6, 10.0, 0.65,
    size=28, bold=True, color=WHITE, font="Trebuchet MS")
txt(s8, "Interactive amortization modeling with visual SHAP-based XAI transparency",
    0.5, 1.25, 10.0, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# image placeholder — main simulator
image_placeholder(s8, 0.5, 1.7, 7.8, 5.5,
    "[ Screenshot: Simulator page with\n  dual-axis line chart (Balance + InterestPaid),\n  extra EMI slider moved, and SHAP panel open ]")

# features right column
feats = [
    ("Dual-Axis Chart",    ACCENT1, "Balance trajectory vs. annual interest paid — rendered via Recharts with live amortization loop."),
    ("Extra EMI Slider",   ACCENT2, "Drag to simulate prepayments. Payoff year updates instantly — all client-side JavaScript math."),
    ("SHAP Analysis Panel",ACCENT3, "Visual XAI showing feature weights:\nIT Stability (+35%), Repo Rate (-22%),\nEMI-Income Ratio (+40%)."),
    ("HITL Trigger",       YELLOW,  "If slider > 2.5× standard EMI, system flags edge case and routes to Human-in-the-Loop."),
]
for i, (ftitle, col, fdesc) in enumerate(feats):
    cy = 1.7 + i * 1.35
    card(s8, 8.6, cy, 4.4, 1.18, fill=CARD_BG, border=col, bw=Pt(1))
    rect(s8, 8.6, cy, 0.07, 1.18, fill_rgb=col)
    txt(s8, ftitle, 8.8, cy+0.1, 4.0, 0.3, size=12, bold=True, color=col,
        font="Trebuchet MS")
    txt(s8, fdesc,  8.8, cy+0.45, 4.0, 0.65, size=11, color=MUTED,
        font="Trebuchet MS")

slide_number(s8, 8)

# ════════════════════════════════════════════════════════════════
#  SLIDE 9 — CSIS & REGULATORY
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_bg(s9)
rect(s9, 0, 0, 13.33, 0.07, fill_rgb=ACCENT3)

tag_box(s9, "DOMAIN-SPECIFIC", 0.5, 0.22, fill=ACCENT3)
txt(s9, "CSIS Subsidy Tracker & Regulatory Action", 0.5, 0.6, 10.0, 0.65,
    size=28, bold=True, color=WHITE, font="Trebuchet MS")
txt(s9, "India-specific financial compliance features built for student borrowers",
    0.5, 1.25, 10.0, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# left: placeholder
image_placeholder(s9, 0.5, 1.7, 6.0, 5.5,
    "[ Screenshot: CSIS Tracker moratorium\n  progress bar + Compliance Alerts panel\n  + Report Unauthorized Loan button ]")

# right: feature boxes
feats9 = [
    ("CSIS Moratorium Tracker", ACCENT3,
     "Tracks 12-month interest subsidy consumption\nagainst moratorium period with a live progress bar."),
    ("Compliance Alert Engine", ACCENT1,
     "Time-coded deadline alerts for:\n• Income proof / EWS certificate renewal\n• NACH mandate validation\n• ITR filing for Section 80E deduction"),
    ("RBI Ombudsman Integration", ACCENT2,
     "One-click 'Report Unauthorized Loan' inside every\nexpanded loan card — simulates regulatory friction\ndirectly with the RBI complaint system."),
    ("Balance Transfer Advisory", YELLOW,
     "AI recommends BoB refinancing at 7.85% vs.\ncurrent 8.5% SBI rate — saving est. ₹2,10,000\nacross the loan tenure for premier graduates."),
]
for i, (ftitle, col, fdesc) in enumerate(feats9):
    cy = 1.7 + i * 1.35
    card(s9, 6.8, cy, 6.1, 1.18, fill=CARD_BG, border=col, bw=Pt(1))
    rect(s9, 6.8, cy, 0.07, 1.18, fill_rgb=col)
    txt(s9, ftitle, 7.0, cy+0.1, 5.7, 0.3, size=12, bold=True, color=col,
        font="Trebuchet MS")
    txt(s9, fdesc,  7.0, cy+0.45, 5.7, 0.65, size=11, color=MUTED,
        font="Trebuchet MS")

slide_number(s9, 9)

# ════════════════════════════════════════════════════════════════
#  SLIDE 10 — LoRA PIPELINE
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_bg(s10)
rect(s10, 0, 0, 13.33, 0.07, fill_rgb=YELLOW)

tag_box(s10, "IP GENERATION", 0.5, 0.22, fill=YELLOW)
txt(s10, "Continuous IP Generation — The LoRA Pipeline", 0.5, 0.6, 10.0, 0.65,
    size=27, bold=True, color=WHITE, font="Trebuchet MS")
txt(s10, "Every verified interaction auto-generates a proprietary fine-tuning dataset entry",
    0.5, 1.25, 10.5, 0.35, size=13, color=MUTED, font="Trebuchet MS")

# left: JSONL screenshot placeholder
image_placeholder(s10, 0.5, 1.7, 6.0, 5.5,
    "[ Screenshot: Code editor showing\n  lora-training-data.jsonl file\n  with formatted JSONL entries ]")

# right: pipeline breakdown
card(s10, 6.8, 1.7, 6.1, 5.5, fill=CARD_BG, border=YELLOW, bw=Pt(1.2))
txt(s10, "How the Pipeline Works", 6.95, 1.85, 5.7, 0.38,
    size=14, bold=True, color=YELLOW, font="Trebuchet MS")

lora_steps = [
    ("TRIGGER",  "confidenceScore ≥ 80 after Critic audit",          ACCENT1),
    ("CAPTURE",  "System prompt + user query + finalAdvice recorded", ACCENT2),
    ("FORMAT",   "Structured as 3-turn JSONL (system/user/assistant)",ACCENT3),
    ("APPEND",   "Logged to lora-training-data.jsonl on backend",     YELLOW),
    ("FUTURE",   "Fine-tune Llama-3 8B — 10× cheaper inference",      GREEN_OK),
]
for i, (step, desc, col) in enumerate(lora_steps):
    cy = 2.42 + i * 0.92
    rect(s10, 7.0, cy, 0.65, 0.32, fill_rgb=col)
    txt(s10, step, 7.0, cy+0.01, 0.65, 0.3, size=8, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER, font="Trebuchet MS")
    txt(s10, desc, 7.8, cy-0.02, 4.9, 0.38, size=11, color=MUTED,
        font="Trebuchet MS")
    if i < 4:
        rect(s10, 7.3, cy+0.32, 0.02, 0.6, fill_rgb=MUTED)

slide_number(s10, 10)

# ════════════════════════════════════════════════════════════════
#  SLIDE 11 — LIVE DEMO
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
add_bg(s11)

# big bg gradient effect via rectangles
rect(s11, 0, 0, 13.33, 7.5, fill_rgb=BG_MID)
rect(s11, 0, 0, 0.12, 7.5, fill_rgb=ACCENT1)

tag_box(s11, "DEMO TIME", 0.5, 0.5, fill=ACCENT1)
txt(s11, "Live Demonstration", 0.5, 1.1, 8.0, 0.75,
    size=36, bold=True, color=WHITE, font="Trebuchet MS")
txt(s11, "Three key scenarios will be demonstrated live",
    0.5, 1.9, 8.0, 0.4, size=15, color=MUTED, font="Trebuchet MS")

demos = [
    ("01", "Dynamic Profile Switching",
     "Switch between Deepak (784 CIBIL), Sumit (650), Kashish (810)\nand Tarun (720) — watch charts and advice update in real-time.", ACCENT1),
    ("02", "Interactive EMI Simulator",
     "Drag the prepayment slider — observe payoff year change\nand SHAP analysis update live with XAI feature weights.", ACCENT2),
    ("03", "The Scam Trap Query",
     "Ask: 'I got a message I won a lottery of ₹5 lakh and they\nwant GST of ₹10,000' — watch the 🚨 Scam Alert fire.", RED_ALERT),
]
for i, (num, title, desc, col) in enumerate(demos):
    cy = 2.6 + i * 1.55
    card(s11, 0.5, cy, 8.5, 1.3, fill=CARD_BG, border=col, bw=Pt(1.5))
    rect(s11, 0.5, cy, 0.08, 1.3, fill_rgb=col)
    # number
    rect(s11, 0.7, cy+0.28, 0.55, 0.55, fill_rgb=col)
    txt(s11, num, 0.7, cy+0.28, 0.55, 0.55, size=14, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER, font="Trebuchet MS")
    txt(s11, title, 1.42, cy+0.15, 7.0, 0.38, size=14, bold=True, color=col,
        font="Trebuchet MS")
    txt(s11, desc,  1.42, cy+0.58, 7.0, 0.62, size=11, color=MUTED,
        font="Trebuchet MS")

# right placeholder
image_placeholder(s11, 9.2, 2.2, 3.9, 4.6,
    "[ Demo placeholder /\n  project running\n  at localhost:5173 ]")

slide_number(s11, 11)

# ════════════════════════════════════════════════════════════════
#  SLIDE 12 — CONCLUSION
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
add_bg(s12, BG_DARK)

# bottom teal band
rect(s12, 0, 6.5, 13.33, 1.0, fill_rgb=RGBColor(0x06,0x1A,0x2C))
rect(s12, 0, 6.5, 13.33, 0.05, fill_rgb=ACCENT1)

# left stripe
rect(s12, 0, 0, 0.12, 7.5, fill_rgb=ACCENT1)

tag_box(s12, "CONCLUSION", 0.5, 0.4, fill=ACCENT1)
txt(s12, "Bridging Generative AI &", 0.5, 0.9, 9.0, 0.6,
    size=30, bold=True, color=WHITE, font="Trebuchet MS")
txt(s12, "Deterministic Financial Engineering", 0.5, 1.5, 9.5, 0.6,
    size=30, bold=True, color=ACCENT1, font="Trebuchet MS")

rect(s12, 0.5, 2.25, 4.5, 0.04, fill_rgb=MUTED)

summary_pts = [
    "Multi-Agent Debating eliminates AI hallucinations & fraud risk",
    "Secure Context Injection ensures full DPDP compliance",
    "Deterministic frontend math provides provably correct outputs",
    "LoRA pipeline generates compounding proprietary IP value",
    "Platform designed for scale: CIBIL API, mobile app, Hindi NLP",
]
for j, pt in enumerate(summary_pts):
    cy = 2.5 + j * 0.7
    rect(s12, 0.5, cy+0.1, 0.2, 0.2, fill_rgb=ACCENT1)
    txt(s12, pt, 0.85, cy+0.03, 6.5, 0.45, size=13, color=MUTED,
        font="Trebuchet MS")

# right: final image placeholder
image_placeholder(s12, 7.3, 0.5, 5.7, 5.7,
    "[ Final Dashboard screenshot\n  or project logo ]")

# Q&A footer
txt(s12, "Thank You  ·  Open for Questions", 0.5, 6.6, 9.0, 0.45,
    size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="Trebuchet MS")
txt(s12, "Deepak Kumar  |  22BCS14872  |  B.E. CSE  |  Chandigarh University",
    0.5, 7.05, 10.0, 0.3, size=10, color=MUTED, font="Trebuchet MS")

slide_number(s12, 12)

# ─── SAVE ────────────────────────────────────────────────────────
file_path = '/Users/deepak/Desktop/Loan_assistant/output.pptx'

prs.save(file_path)

print("Saved successfully.")